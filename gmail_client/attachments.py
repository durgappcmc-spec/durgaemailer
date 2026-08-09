# NOTE: Attachment helpers + document text/context extraction for email drafts.
from __future__ import annotations

import base64
import io
import sys
from typing import Any


def files_to_attachments(uploaded_files: list[Any] | None) -> list[dict[str, Any]]:
    """Convert Streamlit UploadedFile list to send/schedule attachment dicts.

    Any file type is accepted as an email attachment. Text/context is extracted
    when possible (PDF, Office, text, images via Gemini); otherwise a metadata
    stub is still included so drafts know the file is attached.
    """
    out: list[dict[str, Any]] = []
    for f in uploaded_files or []:
        data = f.getvalue()
        name = getattr(f, "name", None) or "file"
        mime = getattr(f, "type", None) or _guess_mime(name)
        item: dict[str, Any] = {
            "name": name,
            "data": data,
            "mime_type": mime,
            "data_base64": base64.b64encode(data).decode("ascii"),
            "mimeType": mime,
            "size": len(data),
        }
        text = extract_file_text(name, data, mime)
        if text:
            item["extracted_text"] = text
            item["has_context"] = True
        else:
            item["extracted_text"] = (
                f"[File attached for email: {name} ({mime}, {len(data)} bytes). "
                "No extractable text — still include/attach this file when sending.]"
            )
            item["has_context"] = False
        out.append(item)
    return out


def _guess_mime(name: str) -> str:
    lower = (name or "").lower()
    mapping = {
        ".pdf": "application/pdf",
        ".txt": "text/plain",
        ".md": "text/markdown",
        ".csv": "text/csv",
        ".json": "application/json",
        ".html": "text/html",
        ".htm": "text/html",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".doc": "application/msword",
        ".xls": "application/vnd.ms-excel",
        ".ppt": "application/vnd.ms-powerpoint",
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif": "image/gif",
        ".webp": "image/webp",
        ".bmp": "image/bmp",
        ".rtf": "application/rtf",
        ".zip": "application/zip",
    }
    for ext, mime in mapping.items():
        if lower.endswith(ext):
            return mime
    return "application/octet-stream"


def extract_file_text(name: str, data: bytes, mime: str = "") -> str:
    """Extract plain text / description from any supported upload for LLM context."""
    lower = (name or "").lower()
    mime = (mime or "").lower()
    try:
        if lower.endswith(".pdf") or "pdf" in mime:
            return _extract_pdf_text(data)
        if lower.endswith(".docx") or "wordprocessingml" in mime:
            return _extract_docx(data)
        if lower.endswith(".xlsx") or "spreadsheetml" in mime:
            return _extract_xlsx(data)
        if lower.endswith(".pptx") or "presentationml" in mime:
            return _extract_pptx(data)
        if lower.endswith((".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp")) or mime.startswith(
            "image/"
        ):
            return _describe_image(name, data, mime or _guess_mime(name))
        if lower.endswith((".txt", ".md", ".csv", ".json", ".html", ".htm", ".log", ".rtf")) or mime.startswith(
            "text/"
        ):
            return data.decode("utf-8", errors="replace")
        # Best-effort UTF-8 for unknown small text-ish blobs
        if len(data) < 200_000 and b"\x00" not in data[:1000]:
            sample = data.decode("utf-8", errors="ignore")
            if sample.strip() and sum(c.isprintable() or c.isspace() for c in sample) / max(
                len(sample), 1
            ) > 0.85:
                return sample
    except Exception as e:
        print(f"[attachments] extract failed for {name}: {e}", file=sys.stderr)
    return ""


def _extract_pdf_text(data: bytes) -> str:
    try:
        from pypdf import PdfReader  # type: ignore
    except Exception as e:
        print(f"[attachments] pypdf not available: {e}", file=sys.stderr)
        return ""
    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for page in reader.pages[:40]:
        try:
            parts.append(page.extract_text() or "")
        except Exception:
            continue
    return "\n".join(parts).strip()[:20000]


def _extract_docx(data: bytes) -> str:
    try:
        from docx import Document  # type: ignore
    except Exception as e:
        print(f"[attachments] python-docx missing: {e}", file=sys.stderr)
        return ""
    doc = Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    for table in doc.tables[:20]:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts).strip()[:20000]


def _extract_xlsx(data: bytes) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore
    except Exception as e:
        print(f"[attachments] openpyxl missing: {e}", file=sys.stderr)
        return ""
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets[:8]:
        parts.append(f"## Sheet: {sheet.title}")
        rows_taken = 0
        for row in sheet.iter_rows(values_only=True):
            vals = ["" if v is None else str(v) for v in row]
            if any(v.strip() for v in vals):
                parts.append("\t".join(vals))
                rows_taken += 1
            if rows_taken >= 80:
                parts.append("…")
                break
    return "\n".join(parts).strip()[:20000]


def _extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:
        print(f"[attachments] python-pptx missing: {e}", file=sys.stderr)
        return ""
    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides[:30], 1):
        bits: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                bits.append(shape.text.strip())
        if bits:
            parts.append(f"## Slide {i}\n" + "\n".join(bits))
    return "\n\n".join(parts).strip()[:20000]


def _describe_image(name: str, data: bytes, mime: str) -> str:
    """Use Gemini to summarize an image for email context."""
    try:
        from core.llm import describe_bytes

        summary = describe_bytes(
            data,
            mime_type=mime or "image/png",
            prompt=(
                f"Describe this image ({name}) for use as context when drafting a "
                "professional outreach email. Include any visible text, logos, "
                "offers, dates, and key facts. Be concise but complete."
            ),
        )
        return (summary or "").strip()[:8000]
    except Exception as e:
        print(f"[attachments] image describe failed for {name}: {e}", file=sys.stderr)
        return f"[Image file {name} attached; description unavailable: {e}]"


def document_context_from_attachments(
    attachments: list[dict[str, Any]] | None,
    *,
    max_chars: int = 18000,
) -> str:
    """Build a labeled context block from every staged attachment."""
    if not attachments:
        return ""
    chunks: list[str] = []
    used = 0
    for att in attachments:
        text = (att.get("extracted_text") or "").strip()
        name = att.get("name") or "document"
        mime = att.get("mime_type") or att.get("mimeType") or ""
        flag = "context" if att.get("has_context", bool(text)) else "attach-only"
        header = f"--- File: {name} ({mime}, {flag}) ---\n"
        if not text:
            text = f"[Attached: {name}]"
        room = max_chars - used - len(header)
        if room <= 0:
            break
        body = text[:room]
        chunks.append(header + body)
        used += len(header) + len(body)
    return "\n\n".join(chunks)
